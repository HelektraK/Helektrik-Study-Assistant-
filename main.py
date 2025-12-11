from dotenv import load_dotenv
load_dotenv()

import os
import json
import uuid
import subprocess
import re
import requests
import numpy as np
from datetime import datetime
from typing import Dict, Any, List, Optional

import fitz  # PyMuPDF
from pptx import Presentation
from fastapi import FastAPI, Request, UploadFile, File, Form, Depends, HTTPException, status
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel, Field

# ==============================================================================
# SECTION 1: DATA MODELS (KG5 - Data Model)
# Pydantic schemas defining how data is organized and validated
# ==============================================================================

class FileEntry(BaseModel):
    """Schema for an uploaded file entry."""
    name: str = Field(..., min_length=1, max_length=255)
    type: str
    text: str = ""
    added_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())

class Session(BaseModel):
    """Schema for a study session."""
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: Optional[str] = None
    created_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    updated_at: str = Field(default_factory=lambda: datetime.utcnow().isoformat())
    files: List[FileEntry] = Field(default_factory=list)

class SessionUpdateRequest(BaseModel):
    """Request schema for updating a session (KG5)."""
    name: Optional[str] = Field(None, min_length=1, max_length=100)

class TranscriptResponse(BaseModel):
    """Response schema for transcript endpoint (KG5, KG7)."""
    text: Optional[str] = None
    name: Optional[str] = None
    added_at: Optional[str] = None
    error: Optional[str] = None

class DeleteResponse(BaseModel):
    """Response schema for delete operations (KG5, KG7)."""
    success: bool
    message: str
    deleted_item: Optional[str] = None

# Validation constants (KG3 - Endpoint Validation)
ALLOWED_EXTENSIONS = {"txt", "pdf", "pptx", "mp4", "m4a", "wav", "webm"}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

def validate_file_extension(filename: str) -> str:
    """Validate file extension (KG3)."""
    if not filename or '.' not in filename:
        raise ValueError("Invalid filename")
    ext = filename.rsplit('.', 1)[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"File type .{ext} not allowed. Allowed: {ALLOWED_EXTENSIONS}")
    return ext

def validate_session_id(session_id: str) -> bool:
    """Validate UUID format (KG3)."""
    pattern = re.compile(r'^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$', re.I)
    return bool(pattern.match(session_id))

# ==============================================================================
# SECTION 2: CONFIGURATION
# ==============================================================================

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(BASE_DIR, "upload")
AUDIO_DIR = os.path.join(UPLOAD_DIR, "audio")
DOC_DIR = os.path.join(UPLOAD_DIR, "docs")
VECTOR_STORE_DIR = os.path.join(UPLOAD_DIR, "vector_stores")
SESSIONS_PATH = os.path.join(BASE_DIR, "sessions.json")

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(AUDIO_DIR, exist_ok=True)
os.makedirs(DOC_DIR, exist_ok=True)
os.makedirs(VECTOR_STORE_DIR, exist_ok=True)

# Ollama Configuration
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "qwen2.5:7b")
OLLAMA_EMBED_MODEL = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")

app = FastAPI(title="Helektron Study Assistant", version="2.0.0")
app.mount("/static", StaticFiles(directory=os.path.join(BASE_DIR, "static")), name="static")
templates = Jinja2Templates(directory=os.path.join(BASE_DIR, "templates"))

# ==============================================================================
# SECTION 3: DATA ACCESS LAYER - SESSION CRUD (KG6 - CRUD Operations)
# Handles persistent data storage
# ==============================================================================

def load_sessions() -> Dict[str, Any]:
    """Load sessions from JSON file (KG6 - Read)."""
    if not os.path.exists(SESSIONS_PATH):
        return {}
    with open(SESSIONS_PATH, "r") as f:
        try:
            return json.load(f)
        except json.JSONDecodeError:
            return {}

def save_sessions(sessions: Dict[str, Any]) -> None:
    """Save sessions to JSON file (KG6 - Persistent Data)."""
    with open(SESSIONS_PATH, "w") as f:
        json.dump(sessions, f, indent=2)

def create_session() -> Session:
    """Create a new session (KG6 - Create)."""
    session = Session()
    sessions = load_sessions()
    sessions[session.id] = session.model_dump()
    save_sessions(sessions)
    return session

def get_session(session_id: str) -> Optional[Session]:
    """Get a session by ID (KG6 - Read)."""
    sessions = load_sessions()
    data = sessions.get(session_id)
    if not data:
        return None
    return Session(**data)

def update_session(session: Session) -> Session:
    """Update a session (KG6 - Update)."""
    sessions = load_sessions()
    session.updated_at = datetime.utcnow().isoformat()
    sessions[session.id] = session.model_dump()
    save_sessions(sessions)
    return session

def delete_session(session_id: str) -> bool:
    """Delete a session (KG6 - Delete)."""
    sessions = load_sessions()
    if session_id not in sessions:
        return False
    del sessions[session_id]
    save_sessions(sessions)
    return True

def get_or_create_session(session_id: Optional[str]) -> Session:
    """Get existing or create new session."""
    if session_id:
        session = get_session(session_id)
        if session:
            return session
    return create_session()

def add_file_to_session(session_id: str, filename: str, file_type: str, text: str) -> Session:
    """Add a file to session (KG6 - Update)."""
    session = get_session(session_id)
    if not session:
        raise ValueError("Session not found")
    session.files.append(FileEntry(name=filename, type=file_type, text=text))
    return update_session(session)

def delete_file_from_session(session_id: str, file_index: int) -> FileEntry:
    """Delete a file from session (KG6 - Delete)."""
    session = get_session(session_id)
    if not session:
        raise ValueError("Session not found")
    if file_index < 0 or file_index >= len(session.files):
        raise IndexError("File index out of range")
    removed = session.files.pop(file_index)
    update_session(session)
    return removed

def get_session_text(session_id: str) -> str:
    """Get combined text from all files in session."""
    session = get_session(session_id)
    if not session:
        return ""
    blocks = [f"--- {f.name} ({f.type}) ---\n{f.text}" for f in session.files if f.text]
    return "\n\n".join(blocks)

# ==============================================================================
# SECTION 4: RAG SERVICE (Vector Store for Retrieval-Augmented Generation)
# ==============================================================================

def get_vector_store_path(session_id: str) -> str:
    """Get path to vector store for a session."""
    return os.path.join(VECTOR_STORE_DIR, f"{session_id}_vs.json")

def load_vector_store(session_id: str) -> List[Dict]:
    """Load vector store for a session."""
    path = get_vector_store_path(session_id)
    if not os.path.exists(path):
        return []
    with open(path, "r") as f:
        try:
            return json.load(f)
        except json.JSONDecodeError:
            return []

def save_vector_store(session_id: str, store: List[Dict]) -> None:
    """Save vector store for a session."""
    path = get_vector_store_path(session_id)
    with open(path, "w") as f:
        json.dump(store, f, indent=2)

def embed_text_ollama(text: str) -> List[float]:
    """Generate embedding using Ollama."""
    try:
        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/embeddings",
            json={"model": OLLAMA_EMBED_MODEL, "prompt": text},
            timeout=60
        )
        response.raise_for_status()
        return response.json()["embedding"]
    except Exception as e:
        raise RuntimeError(f"Ollama embedding failed: {e}")

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks for RAG."""
    if len(text) <= chunk_size:
        return [text] if text.strip() else []
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        
        # Try to break at sentence boundary
        if end < len(text):
            for sep in ['\n\n', '\n', '. ', ' ']:
                last_sep = chunk.rfind(sep)
                if last_sep > chunk_size // 2:
                    chunk = chunk[:last_sep + len(sep)]
                    break
        
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    
    return chunks

def add_to_vector_store(session_id: str, text: str, metadata: Dict = None) -> int:
    """Add document to RAG vector store."""
    if not text.strip():
        return 0
    
    store = load_vector_store(session_id)
    chunks = chunk_text(text)
    
    for chunk in chunks:
        try:
            embedding = embed_text_ollama(chunk)
            store.append({
                "text": chunk,
                "embedding": embedding,
                "metadata": metadata or {}
            })
        except Exception as e:
            print(f"Warning: Failed to embed chunk: {e}")
    
    save_vector_store(session_id, store)
    return len(chunks)

def search_vector_store(session_id: str, query: str, top_k: int = 3) -> List[Dict]:
    """Search RAG vector store using cosine similarity."""
    store = load_vector_store(session_id)
    if not store:
        return []
    
    try:
        query_embedding = embed_text_ollama(query)
    except:
        return []
    
    query_vec = np.array(query_embedding)
    
    results = []
    for item in store:
        doc_vec = np.array(item["embedding"])
        # Cosine similarity
        score = np.dot(query_vec, doc_vec) / (np.linalg.norm(query_vec) * np.linalg.norm(doc_vec))
        results.append({
            "text": item["text"],
            "score": float(score),
            "metadata": item.get("metadata", {})
        })
    
    results.sort(key=lambda x: x["score"], reverse=True)
    return results[:top_k]

def get_rag_context(session_id: str, query: str, top_k: int = 5) -> str:
    """Get relevant context from RAG for a query."""
    results = search_vector_store(session_id, query, top_k)
    if not results:
        return ""
    
    context_parts = []
    for i, r in enumerate(results, 1):
        source = r["metadata"].get("filename", "Unknown")
        context_parts.append(f"[Source {i}: {source}]\n{r['text']}")
    
    return "\n\n---\n\n".join(context_parts)

def clear_vector_store(session_id: str) -> None:
    """Clear RAG vector store for a session."""
    path = get_vector_store_path(session_id)
    if os.path.exists(path):
        os.remove(path)

def rebuild_vector_store(session_id: str) -> None:
    """Rebuild RAG vector store from session files."""
    clear_vector_store(session_id)
    session = get_session(session_id)
    if not session:
        return
    for f in session.files:
        if f.text:
            add_to_vector_store(session_id, f.text, {"filename": f.name, "type": f.type})

# ==============================================================================
# SECTION 5: FILE EXTRACTION SERVICE (KG10 - Separation of Concerns)
# ==============================================================================

def extract_txt(path: str) -> str:
    """Extract text from plain text file."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_pdf(path: str) -> str:
    """Extract text from PDF file."""
    doc = fitz.open(path)
    texts = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(texts)

def extract_pptx(path: str) -> str:
    """Extract text from PowerPoint file."""
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def transcribe_audio(path: str) -> str:
    """Transcribe audio using whisper.cpp."""
    # Find whisper binary
    candidates = [
        os.path.join(BASE_DIR, "whisper.cpp", "build", "bin", "whisper-cli"),
        os.path.join(BASE_DIR, "whisper.cpp", "build", "bin", "main"),
        os.path.join(BASE_DIR, "whisper.cpp", "main"),
    ]
    whisper_bin = next((c for c in candidates if os.path.exists(c)), None)
    if not whisper_bin:
        raise RuntimeError("Whisper binary not found.")

    # Find model
    model_candidates = [
        os.path.join(BASE_DIR, "whisper.cpp", "models", "ggml-base.en.bin"),
        os.path.join(BASE_DIR, "whisper.cpp", "models", "ggml-base.bin"),
        os.path.join(BASE_DIR, "whisper.cpp", "models", "ggml-small.en.bin"),
        os.path.join(BASE_DIR, "whisper.cpp", "models", "ggml-tiny.en.bin"),
    ]
    model_path = next((m for m in model_candidates if os.path.exists(m)), None)
    if not model_path:
        raise RuntimeError("Whisper model not found.")

    # Convert to WAV
    wav_path = path + ".wav"
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", path, "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", wav_path],
            check=True, capture_output=True
        )
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"FFmpeg conversion failed: {e.stderr.decode()}")
    except FileNotFoundError:
        raise RuntimeError("FFmpeg not found.")

    # Run whisper
    out_prefix = path + "_out"
    try:
        subprocess.run(
            [whisper_bin, "-m", model_path, "-f", wav_path, "-otxt", "-of", out_prefix],
            check=True, capture_output=True, text=True
        )
    except subprocess.CalledProcessError:
        try:
            subprocess.run(
                [whisper_bin, "-m", model_path, "-f", wav_path, "--output-txt", "--output-file", out_prefix],
                check=True, capture_output=True, text=True
            )
        except subprocess.CalledProcessError as e:
            raise RuntimeError(f"Whisper failed: {e.stderr}")

    txt_path = out_prefix + ".txt"
    if not os.path.exists(txt_path):
        raise RuntimeError("Transcript file not created.")
    
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        transcript = f.read().strip()

    # Cleanup
    for p in [wav_path, txt_path]:
        if os.path.exists(p):
            try:
                os.remove(p)
            except:
                pass

    return transcript if transcript else "[No speech detected]"

def extract_text_from_file(file_path: str, extension: str) -> str:
    """Extract text from file based on extension."""
    extractors = {
        "txt": extract_txt,
        "pdf": extract_pdf,
        "pptx": extract_pptx,
        "mp4": transcribe_audio,
        "m4a": transcribe_audio,
        "wav": transcribe_audio,
        "webm": transcribe_audio,
    }
    extractor = extractors.get(extension)
    if not extractor:
        raise ValueError(f"Unsupported file type: {extension}")
    return extractor(file_path)

def get_file_category(extension: str) -> str:
    """Get file category from extension."""
    return "audio" if extension in {"mp4", "m4a", "wav", "webm"} else "document"

# ==============================================================================
# SECTION 6: AI SERVICE - OLLAMA (KG10 - Separation of Concerns)
# ==============================================================================

def call_ollama(prompt: str) -> str:
    """Call Ollama LLM model."""
    try:
        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/generate",
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False
            },
            timeout=120
        )
        response.raise_for_status()
        return response.json()["response"].strip()
    except Exception as e:
        return f"[Ollama API Error: {e}]"

def get_summary_prompt(text: str, rag_context: str = "") -> str:
    """Generate summary prompt with optional RAG context."""
    context_section = f"\n---RELEVANT CONTEXT (RAG)---\n{rag_context}\n---END CONTEXT---\n" if rag_context else ""
    return f"""
Based on the following combined study materials, create a structured study summary.
{context_section}
Sections to include:
- **Overview**: Topic and goals
- **Key Concepts**: Major ideas, theories, formulas
- **Detailed Breakdown**: Grouped explanations
- **Important Definitions**: Technical terms
- **Examples/Applications**: Real-world uses
- **Main Takeaways**: What to remember

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

def get_keyterms_prompt(text: str, rag_context: str = "") -> str:
    """Generate key terms prompt."""
    context_section = f"\n---TERMINOLOGY CONTEXT (RAG)---\n{rag_context}\n---END CONTEXT---\n" if rag_context else ""
    return f"""
Extract **10-20 key terms** with definitions from the materials below.
{context_section}
For each term include:
- The term
- 1-2 sentence definition
- Why it matters (optional)

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

def get_questions_prompt(text: str, rag_context: str = "") -> str:
    """Generate questions prompt."""
    context_section = f"\n---KEY CONCEPTS (RAG)---\n{rag_context}\n---END CONTEXT---\n" if rag_context else ""
    return f"""
Create **8-12 practice questions** based on the materials below.
{context_section}
Include:
- Conceptual questions
- Short-answer questions
- Application questions

Do NOT provide answers.

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

def get_resources_prompt(text: str, rag_context: str = "") -> str:
    """Generate resources prompt."""
    context_section = f"\n---TOPICS (RAG)---\n{rag_context}\n---END CONTEXT---\n" if rag_context else ""
    return f"""
Recommend **3-7 external resources** for further study based on the materials below.
{context_section}
For each resource include:
- Title
- Type (textbook, article, video)
- Source
- Why it's relevant

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

# ==============================================================================
# SECTION 7: DEPENDENCY INJECTION (KG4)
# ==============================================================================

def get_current_session(session_id: Optional[str] = None) -> Session:
    """Dependency for getting/creating session (KG4)."""
    return get_or_create_session(session_id)

# ==============================================================================
# SECTION 8: UI ENDPOINTS (KG1, KG8 - UI Endpoints & HTMX)
# ==============================================================================

@app.get("/", response_class=HTMLResponse, tags=["UI"])
def index(request: Request):
    """Main page (KG1 - Endpoint Definition, KG8 - UI Endpoint)."""
    return templates.TemplateResponse("index.html", {"request": request})

@app.post("/upload", response_class=HTMLResponse, tags=["UI"])
async def upload_material(
    request: Request,
    file: UploadFile = File(...),
    session_id: Optional[str] = Form(None),
):
    """
    Upload material (KG1, KG2-POST, KG3-Validation, KG6-Create, KG8-HTMX).
    """
    # Validate file extension (KG3)
    try:
        ext = validate_file_extension(file.filename or "")
    except ValueError as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(e))
    
    # Validate file size (KG3)
    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="File too large (max 50MB)")
    if len(contents) == 0:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="File is empty")
    
    # Validate session_id format (KG3)
    if session_id and not validate_session_id(session_id):
        session_id = None
    
    # Get or create session (KG6 - Create)
    session = get_or_create_session(session_id)
    
    # Save file
    file_type = get_file_category(ext)
    save_dir = AUDIO_DIR if file_type == "audio" else DOC_DIR
    saved_path = os.path.join(save_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(saved_path, "wb") as f:
        f.write(contents)
    
    # Extract text
    try:
        text = extract_text_from_file(saved_path, ext)
    except Exception as e:
        text = f"[Error: {e}]"
    
    # Add to session (KG6 - Create/Update)
    add_file_to_session(session.id, file.filename, file_type, text)
    
    # Add to RAG vector store
    try:
        add_to_vector_store(session.id, text, {"filename": file.filename, "type": file_type})
    except Exception as e:
        print(f"RAG indexing warning: {e}")
    
    # Return HTML fragment (KG8 - HTMX)
    session = get_session(session.id)
    return templates.TemplateResponse(
        "fragments/upload_status.html",
        {"request": request, "session_id": session.id, "files": [f.model_dump() for f in session.files]},
    )

@app.post("/upload_live_audio", response_class=HTMLResponse, tags=["UI"])
async def upload_live_audio(
    request: Request,
    audio: UploadFile = File(...),
    session_id: Optional[str] = Form(None),
):
    """Live audio transcription (KG1, KG2-POST, KG8-HTMX)."""
    # Validate (KG3)
    try:
        ext = validate_file_extension(audio.filename or "recording.webm")
    except ValueError as e:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(e))
    
    contents = await audio.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="File too large")
    
    if session_id and not validate_session_id(session_id):
        session_id = None
    
    session = get_or_create_session(session_id)
    
    filename = audio.filename or "live_recording.webm"
    saved_path = os.path.join(AUDIO_DIR, f"live_{uuid.uuid4().hex}.{ext}")
    with open(saved_path, "wb") as f:
        f.write(contents)
    
    try:
        text = transcribe_audio(saved_path)
    except Exception as e:
        text = f"[Error: {e}]"
    
    add_file_to_session(session.id, filename, "audio-live", text)
    
    # Add to RAG
    try:
        add_to_vector_store(session.id, text, {"filename": filename, "type": "audio-live"})
    except:
        pass
    
    session = get_session(session.id)
    return templates.TemplateResponse(
        "fragments/upload_status.html",
        {"request": request, "session_id": session.id, "files": [f.model_dump() for f in session.files]},
    )

# Study tool endpoints (KG8 - HTMX UI) - All use RAG
@app.get("/summary/{session_id}", response_class=HTMLResponse, tags=["UI", "Study Tools"])
def generate_summary(request: Request, session_id: str):
    """Generate summary with RAG (KG1, KG2-GET, KG3, KG8)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    text = get_session_text(session_id)
    rag_context = get_rag_context(session_id, "main concepts overview summary key ideas")
    content = call_ollama(get_summary_prompt(text, rag_context)) if text.strip() else "No material uploaded yet."
    
    return templates.TemplateResponse("fragments/summary.html", {"request": request, "content": content})

@app.get("/keyterms/{session_id}", response_class=HTMLResponse, tags=["UI", "Study Tools"])
def generate_keyterms(request: Request, session_id: str):
    """Generate key terms with RAG (KG1, KG2-GET, KG3, KG8)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    text = get_session_text(session_id)
    rag_context = get_rag_context(session_id, "definitions terminology vocabulary key terms")
    content = call_ollama(get_keyterms_prompt(text, rag_context)) if text.strip() else "No material uploaded yet."
    
    return templates.TemplateResponse("fragments/keyterms.html", {"request": request, "content": content})

@app.get("/questions/{session_id}", response_class=HTMLResponse, tags=["UI", "Study Tools"])
def generate_questions_view(request: Request, session_id: str):
    """Generate questions with RAG (KG1, KG2-GET, KG3, KG8)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    text = get_session_text(session_id)
    rag_context = get_rag_context(session_id, "important concepts examples problems applications")
    content = call_ollama(get_questions_prompt(text, rag_context)) if text.strip() else "No material uploaded yet."
    
    return templates.TemplateResponse("fragments/questions.html", {"request": request, "content": content})

@app.get("/resources/{session_id}", response_class=HTMLResponse, tags=["UI", "Study Tools"])
def generate_resources_view(request: Request, session_id: str):
    """Generate resources with RAG (KG1, KG2-GET, KG3, KG8)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    text = get_session_text(session_id)
    rag_context = get_rag_context(session_id, "topics subjects domain field area")
    content = call_ollama(get_resources_prompt(text, rag_context)) if text.strip() else "No material uploaded yet."
    
    return templates.TemplateResponse("fragments/resources.html", {"request": request, "content": content})

# Delete file UI endpoint (KG8, KG9 - User Interaction Delete)
@app.delete("/session/{session_id}/file/{file_index}", response_class=HTMLResponse, tags=["UI"])
def delete_file_ui(request: Request, session_id: str, file_index: int):
    """Delete file with HTMX response (KG1, KG2-DELETE, KG6-Delete, KG8, KG9)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    if file_index < 0:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid file index")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    if file_index >= len(session.files):
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")
    
    delete_file_from_session(session_id, file_index)
    rebuild_vector_store(session_id)  # Rebuild RAG after delete
    
    session = get_session(session_id)
    return templates.TemplateResponse(
        "fragments/upload_status.html",
        {"request": request, "session_id": session.id, "files": [f.model_dump() for f in session.files]},
    )

# ==============================================================================
# SECTION 9: API ENDPOINTS (KG1, KG7 - API Endpoints & JSON)
# ==============================================================================

@app.get("/api/transcript/{session_id}", response_model=TranscriptResponse, tags=["API"])
def get_latest_transcript(session_id: str):
    """Get latest transcript as JSON (KG1, KG2-GET, KG7-JSON)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    audio_files = [f for f in session.files if f.type == "audio-live"]
    if not audio_files:
        return TranscriptResponse(error="No transcriptions found")
    
    latest = audio_files[-1]
    return TranscriptResponse(text=latest.text, name=latest.name, added_at=latest.added_at)

@app.get("/api/session/{session_id}", tags=["API"])
def get_session_api(session_id: str):
    """Get session details as JSON (KG1, KG2-GET, KG7-JSON)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    return JSONResponse(status_code=status.HTTP_200_OK, content={
        "id": session.id,
        "name": session.name,
        "created_at": session.created_at,
        "updated_at": session.updated_at,
        "file_count": len(session.files)
    })

@app.put("/api/session/{session_id}", tags=["API"])
def update_session_api(session_id: str, data: SessionUpdateRequest):
    """Update session (KG1, KG2-PUT, KG6-Update, KG7-JSON)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    session = get_session(session_id)
    if not session:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    if data.name is not None:
        session.name = data.name
        update_session(session)
    
    return JSONResponse(status_code=status.HTTP_200_OK, content={
        "id": session.id,
        "name": session.name,
        "message": "Session updated"
    })

@app.delete("/api/session/{session_id}", response_model=DeleteResponse, tags=["API"])
def delete_session_api(session_id: str):
    """Delete session (KG1, KG2-DELETE, KG6-Delete, KG7-JSON)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    if not delete_session(session_id):
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    
    clear_vector_store(session_id)
    return DeleteResponse(success=True, message="Session deleted", deleted_item=session_id)

@app.delete("/api/session/{session_id}/file/{file_index}", response_model=DeleteResponse, tags=["API"])
def delete_file_api(session_id: str, file_index: int):
    """Delete file from session (KG1, KG2-DELETE, KG6-Delete, KG7-JSON)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    if file_index < 0:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid file index")
    
    try:
        deleted = delete_file_from_session(session_id, file_index)
        rebuild_vector_store(session_id)
        return DeleteResponse(success=True, message="File deleted", deleted_item=deleted.name)
    except ValueError:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Session not found")
    except IndexError:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

@app.get("/api/rag/stats/{session_id}", tags=["API", "RAG"])
def get_rag_stats(session_id: str):
    """Get RAG vector store stats (KG7-JSON)."""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid session ID")
    
    store = load_vector_store(session_id)
    sources = set(item.get("metadata", {}).get("filename", "unknown") for item in store)
    
    return JSONResponse(status_code=status.HTTP_200_OK, content={
        "session_id": session_id,
        "chunks": len(store),
        "sources": list(sources)
    })

@app.get("/health", tags=["System"])
def health_check():
    """Health check (KG2 - 200 OK status code)."""
    return JSONResponse(status_code=status.HTTP_200_OK, content={
        "status": "healthy",
        "version": "2.0.0",
        "features": ["RAG", "Whisper", "HTMX", "Ollama"]
    })